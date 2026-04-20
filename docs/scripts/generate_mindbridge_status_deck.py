#!/usr/bin/env python3
"""
Generate MindBridge-Project-Status.pptx — visual, stakeholder-friendly status deck.

Also writes:
  • MindBridge-Project-Status-Google-Slides-Import.pptx — same deck; upload to Google Drive
    and use Open with → Google Slides (Google’s supported import path).
  • docs/google-slides/MindBridge-CreateInGoogleSlides.gs — Google Apps Script you can run
    once to create a native Slides file in your own Drive (plain layouts, same wording).

Requires: pip install python-pptx

Usage (from repo root):
  python3 docs/scripts/generate_mindbridge_status_deck.py
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# MindBridge-inspired palette (from product UI)
TEAL = RGBColor(0x2D, 0x7A, 0x8F)
TEAL_LIGHT = RGBColor(0x4D, 0xB8, 0xA8)
MINT_BG = RGBColor(0xE8, 0xF4, 0xF7)
PAGE_BG = RGBColor(0xF8, 0xFA, 0xFB)
GREEN_OK = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1E, 0x29, 0x3B)


def _blank_slide(prs: Presentation):
    layouts = prs.slide_layouts
    for i in (6, 5, len(layouts) - 1):
        if 0 <= i < len(layouts):
            try:
                return prs.slides.add_slide(layouts[i])
            except Exception:
                continue
    return prs.slides.add_slide(layouts[0])


def _full_bleed_bg(slide, color: RGBColor) -> None:
    sw, sh = prs_slide_size(slide)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def prs_slide_size(slide) -> tuple:
    prs = slide.part.package.presentation_part.presentation
    return prs.slide_width, prs.slide_height


def _add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width: Pt | None = None,
):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_width:
            sh.line.width = line_width
    return sh


def _textbox(slide, left, top, width, height, text: str, *, size=Pt(16), bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _header_bar(slide, title: str, subtitle: str | None = None) -> None:
    w, _ = prs_slide_size(slide)
    bar_h = Inches(1.05)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    _textbox(slide, Inches(0.55), Inches(0.22), w - Inches(1.1), Inches(0.55), title, size=Pt(32), bold=True, color=WHITE)
    if subtitle:
        _textbox(
            slide,
            Inches(0.55),
            Inches(0.68),
            w - Inches(1.1),
            Inches(0.32),
            subtitle,
            size=Pt(14),
            bold=False,
            color=RGBColor(0xD4, 0xE9, 0xF0),
        )


def slide_hero(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    w, h = prs_slide_size(slide)
    _full_bleed_bg(slide, PAGE_BG)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), h - Inches(1.15), w, Inches(1.15))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = TEAL
    bottom.line.fill.background()

    _textbox(
        slide,
        Inches(0.75),
        Inches(1.35),
        w - Inches(1.2),
        Inches(1.2),
        "MindBridge",
        size=Pt(54),
        bold=True,
        color=TEAL,
    )
    _textbox(
        slide,
        Inches(0.75),
        Inches(2.55),
        w - Inches(1.2),
        Inches(0.9),
        "Where we are today",
        size=Pt(36),
        bold=False,
        color=SLATE,
    )
    _textbox(
        slide,
        Inches(0.75),
        Inches(3.45),
        w - Inches(1.2),
        Inches(1.0),
        "March 2026  ·  Abhi · Umer · Abel\nA friendly update for everyone who read our proposal and pitch",
        size=Pt(18),
        bold=False,
        color=SLATE,
    )

    # Decorative “card” hint
    card = _add_rect(slide, Inches(0.75), Inches(4.85), Inches(4.2), Inches(1.35), MINT_BG, TEAL, Pt(1))
    tf = card.text_frame
    tf.clear()
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "Mood · Resources · Counselling"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p2 = tf.add_paragraph()
    p2.text = "One place to check in and get support"
    p2.font.size = Pt(13)
    p2.font.color.rgb = SLATE

    _textbox(
        slide,
        Inches(0.75),
        h - Inches(0.85),
        w - Inches(1.0),
        Inches(0.45),
        "Project status deck",
        size=Pt(14),
        bold=False,
        color=WHITE,
    )


def slide_three_pillars(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _header_bar(slide, "At a glance", "Three things MindBridge is built around")
    w, _ = prs_slide_size(slide)
    gap = Inches(0.35)
    card_w = (w - Inches(1.1) - 2 * gap) / 3
    left0 = Inches(0.55)
    top = Inches(1.45)
    card_h = Inches(3.2)

    pillars = [
        ("Mood care", "Quick check-ins, notes, and a simple view of how you’ve been feeling."),
        ("Wellness resources", "Articles, exercises, and tools in one calm library."),
        ("Human support", "See counsellors, book a time, and keep track of sessions."),
    ]
    colors = [TEAL_LIGHT, TEAL, RGBColor(0x3A, 0x8F, 0xA3)]
    for i, ((title, body), c) in enumerate(zip(pillars, colors)):
        x = left0 + i * (card_w + gap)
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, card_w, Inches(0.22))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = c
        top_bar.line.fill.background()
        card = _add_rect(slide, x, top + Inches(0.18), card_w, card_h - Inches(0.18), MINT_BG, c, Pt(1.5))
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.35)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(15)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(10)

    # Progress hint
    _textbox(slide, Inches(0.55), Inches(5.0), w - Inches(1.1), Inches(0.35), "How much of the full vision is live today?", size=Pt(14), bold=True, color=SLATE)
    bar_x, bar_y, bar_w = Inches(0.55), Inches(5.45), w - Inches(1.1)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, bar_w, Inches(0.38))
    bg.adjustments[0] = 0.5
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    bg.line.fill.background()
    fill_w = bar_w * 0.72
    fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, bar_y, fill_w, Inches(0.38))
    fg.adjustments[0] = 0.5
    fg.fill.solid()
    fg.fill.fore_color.rgb = TEAL_LIGHT
    fg.line.fill.background()
    _textbox(slide, bar_x + bar_w - Inches(1.4), bar_y - Inches(0.28), Inches(1.4), Inches(0.25), "~70%", size=Pt(13), bold=True, color=TEAL)


def slide_live_today(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _header_bar(slide, "Live in the app", "What already works for users today")
    w, _ = prs_slide_size(slide)
    items = [
        "Secure sign-up and sign-in; patients and counsellors each see what fits their role.",
        "Mood logging with notes, history, and a simple dashboard view.",
        "A wellness resource hub—content in one organized place.",
        "Browse counsellors, pick a slot, book a session, see upcoming and past visits.",
        "Counsellors can review sessions and mark progress (e.g. accepted, completed).",
    ]
    y = Inches(1.35)
    for line in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.08), Inches(0.28), Inches(0.28))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN_OK
        dot.line.fill.background()
        dtf = dot.text_frame
        dtf.clear()
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        dp = dtf.paragraphs[0]
        dp.text = "✓"
        dp.font.size = Pt(13)
        dp.font.bold = True
        dp.font.color.rgb = WHITE
        dp.alignment = PP_ALIGN.CENTER
        row = _add_rect(slide, Inches(1.05), y, w - Inches(1.45), Inches(0.85), RGBColor(0xF0, 0xFD, 0xF4), GREEN_OK, Pt(0.75))
        tf = row.text_frame
        tf.clear()
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += Inches(0.98)


def slide_still_building(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _header_bar(slide, "Still on the way", "Honest gaps — no jargon")
    w, _ = prs_slide_size(slide)
    items = [
        "Chat button opens a panel, but the live assistant isn’t hooked up there yet.",
        "Community and profile areas show “coming soon.”",
        "No automatic reminders by email or text before appointments.",
        "Patients can’t reschedule or cancel inside the app yet; no built-in video links.",
        "We still don’t diagnose, prescribe, handle insurance, or replace emergency services — on purpose.",
    ]
    y = Inches(1.35)
    for line in items:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, Inches(0.12), Inches(0.88))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = AMBER
        stripe.line.fill.background()
        row = _add_rect(slide, Inches(0.75), y, w - Inches(1.15), Inches(0.88), RGBColor(0xFF, 0xFB, 0xEB), AMBER, Pt(0.75))
        tf = row.text_frame
        tf.clear()
        tf.margin_left = Inches(0.22)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        y += Inches(1.0)


def slide_split_promise(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, PAGE_BG)
    _header_bar(slide, "Compared to our early plan", "What we said vs. what you see now")
    w, h = prs_slide_size(slide)
    mid = Inches(0.4)
    col_w = (w - Inches(1.1) - mid) / 2
    x1 = Inches(0.55)
    x2 = x1 + col_w + mid
    top = Inches(1.35)
    col_h = h - top - Inches(0.65)

    left = _add_rect(slide, x1, top, col_w, col_h, WHITE, TEAL_LIGHT, Pt(1.5))
    tf = left.text_frame
    tf.clear()
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "Mostly matches"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEAL
    for t in [
        "The core idea is alive: mood care, resources, and real counsellor booking.",
        "Safety boundaries are unchanged—we connect and guide, not replace crisis or clinical care.",
    ]:
        p2 = tf.add_paragraph()
        p2.text = t
        p2.font.size = Pt(16)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(14)
        p2.level = 0

    right = _add_rect(slide, x2, top, col_w, col_h, WHITE, AMBER, Pt(1.5))
    tf2 = right.text_frame
    tf2.clear()
    tf2.margin_left = Inches(0.25)
    tf2.margin_right = Inches(0.25)
    tf2.margin_top = Inches(0.35)
    p = tf2.paragraphs[0]
    p.text = "Needs more time"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = AMBER
    for t in [
        "Groups, deep moderation, and rich analytics for schools or clinics.",
        "Extras like saved chat history everywhere we described.",
    ]:
        p2 = tf2.add_paragraph()
        p2.text = t
        p2.font.size = Pt(16)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(14)


def slide_timeline(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _header_bar(slide, "How the semester-style plan lines up", "Four phases, plain language")
    w, _ = prs_slide_size(slide)
    steps = [
        ("Discover & plan", "Research, privacy, and requirements — captured in our proposal."),
        ("Build the foundation", "Accounts, screens, and main journeys — largely what you can click today."),
        ("Grow the features", "Mood, resources, and scheduling feel strong; assistant & community still catching up."),
        ("Polish & ship", "Testing, accessibility, stable demo — next focus."),
    ]
    n = len(steps)
    pad = Inches(0.5)
    step_w = (w - pad * 2 - Inches(0.35) * (n - 1)) / n
    y = Inches(1.85)
    for i, (title, desc) in enumerate(steps):
        x = pad + i * (step_w + Inches(0.35))
        sh = _add_rect(slide, x, y, step_w, Inches(2.55), MINT_BG if i < 3 else RGBColor(0xFE, 0xF3, 0xC7), TEAL, Pt(1))
        tf = sh.text_frame
        tf.clear()
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.2)
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), y + Inches(0.12), Inches(0.42), Inches(0.42))
        num.fill.solid()
        num.fill.fore_color.rgb = TEAL if i < 3 else AMBER
        num.line.fill.background()
        nft = num.text_frame
        nft.clear()
        np = nft.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER
        nft.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = TEAL
        p.space_before = Pt(36)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = SLATE
        p2.space_before = Pt(8)
        if i < n - 1:
            ax = x + step_w + Inches(0.02)
            _textbox(
                slide,
                ax,
                y + Inches(1.0),
                Inches(0.32),
                Inches(0.5),
                "→",
                size=Pt(28),
                bold=True,
                color=TEAL_LIGHT,
                align=PP_ALIGN.CENTER,
            )


def slide_next(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _full_bleed_bg(slide, WHITE)
    _header_bar(slide, "What we’re tackling next", "Near-term priorities")
    w, _ = prs_slide_size(slide)
    items = [
        "Connect the in-app chat so answers feel real and helpful.",
        "Let people change or cancel bookings, then add gentle reminders.",
        "Choose: a small community slice, or double down on polish for demos.",
        "Smooth the whole experience: wording, flow, and reliability for real users.",
    ]
    y = Inches(1.4)
    for i, line in enumerate(items, start=1):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        bt = badge.text_frame
        bt.clear()
        bp = bt.paragraphs[0]
        bp.text = str(i)
        bp.font.size = Pt(22)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        row = _add_rect(slide, Inches(1.35), y - Inches(0.02), w - Inches(1.75), Inches(0.78), MINT_BG, TEAL, Pt(0.75))
        tf = row.text_frame
        tf.clear()
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.18)
        p = tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        y += Inches(1.05)

    _textbox(
        slide,
        Inches(0.55),
        Inches(6.35),
        w - Inches(1.1),
        Inches(0.55),
        "Long-form detail: Group 4 Project Proposal · MindBridge pitch deck",
        size=Pt(13),
        color=SLATE,
    )


def _write_google_apps_script(path: Path) -> None:
    """Emit SlidesApp script: run at script.google.com → creates a new presentation in Drive."""

    def j(s: str) -> str:
        return json.dumps(s, ensure_ascii=False)

    title_main = "MindBridge — Project Status (Mar 2026)"
    hero_sub = (
        "Where we are today\n"
        "March 2026 · Abhi · Umer · Abel\n"
        "A friendly update for everyone who read our proposal and pitch\n\n"
        "Mood · Resources · Counselling — one place to check in and get support"
    )
    pillars_body = (
        "Three things MindBridge is built around:\n\n"
        "• Mood care — Quick check-ins, notes, and a simple view of how you’ve been feeling.\n"
        "• Wellness resources — Articles, exercises, and tools in one calm library.\n"
        "• Human support — See counsellors, book a time, and keep track of sessions.\n\n"
        "Roughly 70% of the full vision is live in the app today."
    )
    live_body = (
        "What already works for users today:\n\n"
        "• Secure sign-up and sign-in; patients and counsellors each see what fits their role.\n"
        "• Mood logging with notes, history, and a simple dashboard view.\n"
        "• A wellness resource hub—content in one organized place.\n"
        "• Browse counsellors, pick a slot, book a session, see upcoming and past visits.\n"
        "• Counsellors can review sessions and mark progress (e.g. accepted, completed)."
    )
    gaps_body = (
        "Honest gaps — no jargon:\n\n"
        "• Chat opens a panel, but the live assistant isn’t hooked up there yet.\n"
        "• Community and profile areas show “coming soon.”\n"
        "• No automatic reminders by email or text before appointments.\n"
        "• Patients can’t reschedule or cancel inside the app yet; no built-in video links.\n"
        "• We don’t diagnose, prescribe, handle insurance, or replace emergency services — on purpose."
    )
    compare_body = (
        "Mostly matches\n"
        "• The core idea is alive: mood care, resources, and real counsellor booking.\n"
        "• Safety boundaries are unchanged—we connect and guide, not replace crisis or clinical care.\n\n"
        "Needs more time\n"
        "• Groups, deep moderation, and rich analytics for schools or clinics.\n"
        "• Extras like saved chat history everywhere we described."
    )
    timeline_body = (
        "Four phases, plain language:\n\n"
        "1. Discover & plan — Research, privacy, and requirements (in our proposal).\n"
        "2. Build the foundation — Accounts, screens, and main journeys (largely what you can click today).\n"
        "3. Grow the features — Mood, resources, and scheduling feel strong; assistant & community still catching up.\n"
        "4. Polish & ship — Testing, accessibility, stable demo — next focus."
    )
    next_body = (
        "Near-term priorities:\n\n"
        "1. Connect the in-app chat so answers feel real and helpful.\n"
        "2. Let people change or cancel bookings, then add gentle reminders.\n"
        "3. Choose: a small community slice, or double down on polish for demos.\n"
        "4. Smooth the whole experience: wording, flow, and reliability for real users.\n\n"
        "Long-form detail: Group 4 Project Proposal · MindBridge pitch deck."
    )
    thanks_sub = "Questions & feedback welcome ♥"

    slides_js = ",\n    ".join(
        [
            f"{{ layout: 'TITLE', title: {j('MindBridge')}, body: {j(hero_sub)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('At a glance')}, body: {j(pillars_body)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('Live in the app')}, body: {j(live_body)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('Still on the way')}, body: {j(gaps_body)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('Compared to our early plan')}, body: {j(compare_body)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('How the semester-style plan lines up')}, body: {j(timeline_body)} }}",
            f"{{ layout: 'TITLE_AND_BODY', title: {j('What we’re tackling next')}, body: {j(next_body)} }}",
            f"{{ layout: 'TITLE', title: {j('Thank you')}, body: {j(thanks_sub)} }}",
        ]
    )

    src = f"""/* MindBridge — create a native Google Slides deck in your Drive
 *
 * Google does not use a downloadable “.gslides” file on your computer. You can either:
 *   A) Upload MindBridge-Project-Status-Google-Slides-Import.pptx to Drive → Open with → Google Slides, or
 *   B) Run this script once (same wording; simpler layouts than the PowerPoint).
 *
 * Steps:
 *   1. Open https://script.google.com and sign in.
 *   2. New project → paste this file → Save.
 *   3. Select function “createMindBridgeProjectStatusDeck” → Run → Allow access.
 *   4. Open the new file from the link in the execution log, or find “{title_main}” in Google Drive.
 */

function createMindBridgeProjectStatusDeck() {{
  var title = {j(title_main)};
  var pres = SlidesApp.create(title);
  pres.appendSlide(SlidesApp.PredefinedLayout.TITLE);
  pres.getSlides()[0].remove();

  var slides = [
    {slides_js}
  ];

  var L = SlidesApp.PredefinedLayout;
  for (var i = 0; i < slides.length; i++) {{
    var spec = slides[i];
    var slide;
    if (i === 0) {{
      slide = pres.getSlides()[0];
    }} else {{
      slide = pres.appendSlide(spec.layout === 'TITLE' ? L.TITLE : L.TITLE_AND_BODY);
    }}
    var placeholders = slide.getPlaceholders();
    if (spec.layout === 'TITLE') {{
      if (placeholders.length >= 1)
        placeholders[0].asShape().getText().setText(spec.title);
      if (placeholders.length >= 2)
        placeholders[1].asShape().getText().setText(spec.body);
    }} else {{
      if (placeholders.length >= 1)
        placeholders[0].asShape().getText().setText(spec.title);
      if (placeholders.length >= 2)
        placeholders[1].asShape().getText().setText(spec.body);
    }}
  }}
  var url = pres.getUrl();
  Logger.log('Created: ' + url);
}}
"""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(src, encoding="utf-8")


def slide_closing(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    w, h = prs_slide_size(slide)
    _full_bleed_bg(slide, TEAL)
    _textbox(
        slide,
        Inches(0.75),
        Inches(2.2),
        w - Inches(1.5),
        Inches(1.2),
        "Thank you",
        size=Pt(48),
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(0.75),
        Inches(3.5),
        w - Inches(1.5),
        Inches(0.8),
        "Questions & feedback welcome",
        size=Pt(24),
        bold=False,
        color=RGBColor(0xD4, 0xE9, 0xF0),
        align=PP_ALIGN.CENTER,
    )
    heart = slide.shapes.add_shape(MSO_SHAPE.HEART, w / 2 - Inches(0.35), Inches(4.35), Inches(0.7), Inches(0.65))
    heart.fill.solid()
    heart.fill.fore_color.rgb = TEAL_LIGHT
    heart.line.fill.background()


def main() -> None:
    root = Path(__file__).resolve().parents[2]
    out = root / "MindBridge-Project-Status.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_hero(prs)
    slide_three_pillars(prs)
    slide_live_today(prs)
    slide_still_building(prs)
    slide_split_promise(prs)
    slide_timeline(prs)
    slide_next(prs)
    try:
        slide_closing(prs)
    except Exception:
        # Older python-pptx may lack HEART shape on some builds
        slide = _blank_slide(prs)
        w, h = prs_slide_size(slide)
        _full_bleed_bg(slide, TEAL)
        _textbox(
            slide,
            Inches(0.75),
            Inches(2.5),
            w - Inches(1.5),
            Inches(1.5),
            "Thank you\n\nQuestions welcome ♥",
            size=Pt(40),
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    prs.save(out)

    import_copy = root / "MindBridge-Project-Status-Google-Slides-Import.pptx"
    shutil.copy2(out, import_copy)

    gs_path = root / "docs" / "google-slides" / "MindBridge-CreateInGoogleSlides.gs"
    _write_google_apps_script(gs_path)

    print(f"Wrote {out}")
    print(f"Wrote {import_copy}  (Google Drive → Upload → Open with → Google Slides)")
    print(f"Wrote {gs_path}  (script.google.com → paste → Run → native Slides in your Drive)")


if __name__ == "__main__":
    main()
